from typing import List, Dict, Any
import logging
import os

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".pdf", ".pptx", ".txt", ".md", ".docx",
    ".mp4", ".mkv", ".mov", ".avi", ".webm",
    ".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac"
}

TEXT_EXTENSIONS = {".txt", ".md"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac"}
VIDEO_EXTENSIONS = {".mp4", ".mkv", ".mov", ".avi", ".webm"}

def parse_file(path: str) -> List[Dict[str, Any]]:
    """Dispatch to the right parser based on extension."""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return parse_pdf(path)
    elif ext == ".pptx":
        return parse_pptx(path)
    elif ext in TEXT_EXTENSIONS:
        return parse_text(path)
    elif ext == ".docx":
        return parse_docx(path)
    elif ext in AUDIO_EXTENSIONS:
        return parse_audio(path)
    elif ext in VIDEO_EXTENSIONS:
        return parse_video(path)
    else:
        raise ValueError(f"Unsupported file extension: {ext}")

def parse_text(path: str) -> List[Dict[str, Any]]:
    """Plain text or Markdown file extraction."""
    logger.info(f"Parsing text/markdown: {path}")
    source = os.path.basename(path)
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()

    if not content.strip():
        return []

    # If text is long, break it into pseudo-pages of ~2000 characters
    page_size = 2000
    pages = []
    lines = content.splitlines(keepends=True)
    current_page_text = []
    current_len = 0
    page_num = 1

    for line in lines:
        current_page_text.append(line)
        current_len += len(line)
        if current_len >= page_size:
            text = "".join(current_page_text).strip()
            if text:
                pages.append({"text": text, "page": page_num, "source": source})
                page_num += 1
            current_page_text = []
            current_len = 0

    if current_page_text:
        text = "".join(current_page_text).strip()
        if text:
            pages.append({"text": text, "page": page_num, "source": source})

    return pages

def parse_docx(path: str) -> List[Dict[str, Any]]:
    """Extract text from Word .docx file using python-docx or zipfile XML fallback."""
    logger.info(f"Parsing DOCX: {path}")
    source = os.path.basename(path)
    try:
        import docx
        doc = docx.Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        full_text = "\n\n".join(paragraphs)
    except Exception:
        # Fallback to direct XML extraction from docx zip
        import zipfile
        import xml.etree.ElementTree as ET
        full_text = ""
        try:
            with zipfile.ZipFile(path) as z:
                xml_content = z.read("word/document.xml")
                tree = ET.fromstring(xml_content)
                text_elements = [elem.text for elem in tree.iter() if elem.text]
                full_text = " ".join(text_elements)
        except Exception as exc:
            logger.error(f"Failed to parse docx {path}: {exc}")

    if not full_text.strip():
        return []

    # Break into pseudo-pages
    page_size = 2000
    pages = []
    chunks = [full_text[i:i + page_size] for i in range(0, len(full_text), page_size)]
    for i, c in enumerate(chunks, start=1):
        if c.strip():
            pages.append({"text": c.strip(), "page": i, "source": source})
    return pages

def parse_pdf(path: str) -> List[Dict[str, Any]]:
    """PyMuPDF page-by-page extraction."""
    import fitz
    logger.info(f"Parsing PDF: {path}")
    doc = fitz.open(path)
    pages = []
    source = os.path.basename(path)
    for i in range(len(doc)):
        page = doc.load_page(i)
        text = page.get_text()
        if text.strip():
            pages.append({
                "text": text.strip(),
                "page": i + 1,
                "source": source
            })
    return pages

def parse_pptx(path: str) -> List[Dict[str, Any]]:
    """python-pptx slide-by-slide text extraction."""
    from pptx import Presentation
    logger.info(f"Parsing PPTX: {path}")
    prs = Presentation(path)
    pages = []
    source = os.path.basename(path)
    for i, slide in enumerate(prs.slides):
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
        slide_text = "\n".join(text).strip()
        if slide_text:
            pages.append({
                "text": slide_text,
                "page": i + 1,
                "source": source
            })
    return pages

def _transcribe_with_whisper(audio_path: str, source: str) -> List[Dict[str, Any]]:
    """Faster-Whisper transcription. Returns segments grouped as pages."""
    from faster_whisper import WhisperModel
    from app.config import get_settings
    settings = get_settings()
    logger.info(f"Transcribing {audio_path} using Whisper ({settings.whisper_model})")
    
    model = WhisperModel(settings.whisper_model, device=settings.whisper_device, compute_type="int8")
    segments, _ = model.transcribe(audio_path, beam_size=5)
    
    pages = []
    current_text = []
    current_page = 1
    # Group every ~30 seconds of speech into one "page"
    page_duration = 30.0
    start_time = 0.0

    for segment in segments:
        current_text.append(segment.text)
        if segment.end - start_time >= page_duration:
            pages.append({
                "text": " ".join(current_text).strip(),
                "page": current_page,
                "source": source
            })
            current_text = []
            current_page += 1
            start_time = segment.end

    if current_text:
        pages.append({
            "text": " ".join(current_text).strip(),
            "page": current_page,
            "source": source
        })

    return pages

def parse_audio(path: str) -> List[Dict[str, Any]]:
    source = os.path.basename(path)
    return _transcribe_with_whisper(path, source)

def parse_video(path: str) -> List[Dict[str, Any]]:
    """Extract audio with ffmpeg then transcribe."""
    import subprocess, tempfile
    source = os.path.basename(path)
    logger.info(f"Extracting audio from video: {path}")
    
    with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
        audio_path = tmp.name
        
    try:
        subprocess.run(
            ["ffmpeg", "-y", "-i", path, "-vn", "-acodec", "pcm_s16le", "-ar", "16000", "-ac", "1", audio_path],
            check=True, capture_output=True
        )
        return _transcribe_with_whisper(audio_path, source)
    finally:
        if os.path.exists(audio_path):
            os.unlink(audio_path)
